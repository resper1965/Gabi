"""
Presentation Generator — Creates PowerPoint slides from legal document content.

Uses Gemini to extract key points, then renders them via python-pptx
into professional-looking slides.
"""

import io
import json
import logging
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.core.ai import generate, safe_parse_json

logger = logging.getLogger("gabi.presentation")

SLIDE_PROMPT = """Analise o conteúdo abaixo e crie uma estrutura de apresentação profissional.

RETORNE EXCLUSIVAMENTE JSON válido:
{{
  "title": "Título da apresentação",
  "subtitle": "Subtítulo ou contexto",
  "slides": [
    {{
      "title": "Título do slide",
      "bullets": ["Ponto 1", "Ponto 2", "Ponto 3"],
      "notes": "Notas do apresentador (opcional)"
    }}
  ],
  "conclusion": {{
    "title": "Conclusão",
    "bullets": ["Ponto final 1", "Ponto final 2"]
  }}
}}

Regras:
- Máximo 10 slides (sem contar capa e conclusão)
- 3-5 bullets por slide
- Bullets concisos (máx 15 palavras cada)
- Linguagem formal e objetiva
- Foque nos pontos mais relevantes para um advogado/executivo

CONTEÚDO:
{content}"""

# Color palette — dark professional theme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)


async def generate_presentation(
    content: str,
    theme: str = "professional",
    custom_title: Optional[str] = None,
) -> bytes:
    """
    Generate a PowerPoint presentation from text content.

    Args:
        content: Document text to base the presentation on.
        theme: Visual theme (currently only 'professional').
        custom_title: Override the AI-generated title.

    Returns:
        Bytes of the .pptx file.
    """
    # Step 1: Extract structure via Gemini
    snippet = content[:8000]
    prompt = SLIDE_PROMPT.format(content=snippet)

    try:
        raw = await generate(module="flash", prompt=prompt)
        structure = safe_parse_json(raw)
    except Exception as e:
        logger.error("Slide structure generation failed: %s", e)
        structure = {
            "title": custom_title or "Apresentação",
            "subtitle": "",
            "slides": [{"title": "Conteúdo", "bullets": ["Erro ao gerar estrutura de slides"]}],
            "conclusion": {"title": "Conclusão", "bullets": []},
        }

    if custom_title:
        structure["title"] = custom_title

    # Step 2: Render via python-pptx
    pptx_bytes = _render_pptx(structure)
    return pptx_bytes


def _render_pptx(structure: dict) -> bytes:
    """Render the slide structure into a PowerPoint file."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title Slide
    _add_title_slide(prs, structure.get("title", ""), structure.get("subtitle", ""))

    # Content Slides
    for slide_data in structure.get("slides", []):
        _add_content_slide(
            prs,
            slide_data.get("title", ""),
            slide_data.get("bullets", []),
            slide_data.get("notes", ""),
        )

    # Conclusion Slide
    conclusion = structure.get("conclusion", {})
    if conclusion:
        _add_content_slide(
            prs,
            conclusion.get("title", "Conclusão"),
            conclusion.get("bullets", []),
            is_conclusion=True,
        )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _set_slide_bg(slide, color: RGBColor):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(prs: Presentation, title: str, subtitle: str):
    """Create a dark-themed title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, DARK_BG)

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(11)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    top2 = Inches(4.2)
    txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.LEFT

    # Accent line
    from pptx.shapes.autoshape import Shape
    left_line = Inches(1)
    top_line = Inches(4.0)
    width_line = Inches(3)
    height_line = Emu(36000)  # ~1pt
    shape = slide.shapes.add_shape(1, left_line, top_line, width_line, height_line)  # Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    notes: str = "",
    is_conclusion: bool = False,
):
    """Create a content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _set_slide_bg(slide, DARK_BG)

    # Slide title
    left = Inches(1)
    top = Inches(0.8)
    width = Inches(11)
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28) if not is_conclusion else Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE if is_conclusion else WHITE

    # Accent line under title
    shape = slide.shapes.add_shape(1, left, Inches(1.7), Inches(2), Emu(27000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # Bullets
    bullet_top = Inches(2.2)
    txBox2 = slide.shapes.add_textbox(Inches(1.2), bullet_top, Inches(10.5), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()

        p.text = f"▸  {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

    # Speaker notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
